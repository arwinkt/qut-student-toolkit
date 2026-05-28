"""
scan_school.py — School Folder → Obsidian Notes
================================================
Scans your school folder (default: ~/Desktop/school) and converts ALL useful
file types into readable markdown notes in second_brain/education/

SUPPORTED:
  Documents : .pdf, .docx, .doc, .txt, .md
  Slides    : .pptx, .pptm, .ppsx, .ppt
  Spreadsheets: .xlsx, .xls, .csv
  Images    : .png, .jpg, .jpeg, .gif, .webp, .bmp (described via Claude Vision)
  Diagrams  : .drawio (XML extracted + summarised)
  Code/Data : .py, .js, .html, .css, .json, .xml, .yaml, .yml

INSTALL DEPENDENCIES (run once):
    pip install pdfplumber python-docx python-pptx openpyxl anthropic Pillow

THEN RUN:
    python scripts/scan_school.py

Filter to one subject:
    python scripts/scan_school.py --subject "Database Management"

Skip image analysis (faster, no API calls):
    python scripts/scan_school.py --no-vision
"""

import os
import re
import sys
import json
import base64
import argparse
import xml.etree.ElementTree as ET
from pathlib import Path
from datetime import datetime

# ── Paths ──────────────────────────────────────────────────────────────────
SCHOOL_DIR = Path(os.environ.get("SCHOOL_DIR") or Path.home() / "Desktop" / "school")  # Override with SCHOOL_DIR env var
VAULT_DIR  = Path(os.environ.get("VAULT_DIR") or Path(__file__).parent.parent)  # Override with VAULT_DIR env var
OUTPUT_DIR = VAULT_DIR / "education"

# ── File type groups ───────────────────────────────────────────────────────
DOCUMENT_EXTS  = {".pdf", ".docx", ".doc", ".txt", ".md"}
SLIDE_EXTS     = {".pptx", ".pptm", ".ppsx", ".ppt"}
SHEET_EXTS     = {".xlsx", ".xls", ".csv"}
IMAGE_EXTS     = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
DIAGRAM_EXTS   = {".drawio"}
CODE_EXTS      = {".py", ".js", ".html", ".css", ".json", ".xml", ".yaml", ".yml"}
ALL_SUPPORTED  = DOCUMENT_EXTS | SLIDE_EXTS | SHEET_EXTS | IMAGE_EXTS | DIAGRAM_EXTS | CODE_EXTS

SKIP_DIRS  = {".obsidian", ".smart-env", ".claude", "node_modules", "__pycache__"}
SKIP_FILES = {"thumbs.db", ".ds_store"}

# ── Utility ────────────────────────────────────────────────────────────────

def safe_filename(name: str) -> str:
    name = re.sub(r'[\\/:*?"<>|]', "", name)
    return re.sub(r"\s+", " ", name).strip()


def file_to_base64(path: Path) -> str:
    with open(path, "rb") as f:
        return base64.standard_b64encode(f.read()).decode("utf-8")


# ── Extractors ─────────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    parts.append(f"<!-- Page {i} -->\n{text.strip()}")
        return "\n\n".join(parts) or "[No extractable text found in PDF]"
    except ImportError:
        return "[ERROR: pdfplumber not installed — pip install pdfplumber]"
    except Exception as e:
        return f"[ERROR reading PDF: {e}]"


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            style = para.style.name if para.style else ""
            if   "Heading 1" in style: parts.append(f"# {para.text.strip()}")
            elif "Heading 2" in style: parts.append(f"## {para.text.strip()}")
            elif "Heading 3" in style: parts.append(f"### {para.text.strip()}")
            else:                      parts.append(para.text.strip())
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                rows.insert(1, sep)
                parts.append("\n".join(rows))
        return "\n\n".join(parts) or "[No extractable text found in DOCX]"
    except ImportError:
        return "[ERROR: python-docx not installed — pip install python-docx]"
    except Exception as e:
        return f"[ERROR reading DOCX: {e}]"


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            # Slide title first if available
            if slide.shapes.title and slide.shapes.title.text.strip():
                texts.append(f"**{slide.shapes.title.text.strip()}**")
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                # Extract table content from slides
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                    if rows:
                        sep = "| " + " | ".join(["---"] * len(shape.table.rows[0].cells)) + " |"
                        rows.insert(1, sep)
                        texts.append("\n".join(rows))
            if texts:
                slides.append(f"### Slide {i}\n\n" + "\n\n".join(texts))
        return "\n\n---\n\n".join(slides) or "[No extractable text found in PPTX]"
    except ImportError:
        return "[ERROR: python-pptx not installed — pip install python-pptx]"
    except Exception as e:
        return f"[ERROR reading PPTX: {e}]"


def extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            md_rows = []
            for row in rows:
                if any(c is not None for c in row):
                    md_rows.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
            if md_rows:
                sep = "| " + " | ".join(["---"] * len(rows[0])) + " |"
                md_rows.insert(1, sep)
                parts.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(md_rows))
        return "\n\n".join(parts) or "[No data found in XLSX]"
    except ImportError:
        return "[ERROR: openpyxl not installed — pip install openpyxl]"
    except Exception as e:
        return f"[ERROR reading XLSX: {e}]"


def extract_csv(path: Path) -> str:
    try:
        import csv
        rows = []
        with open(path, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append("| " + " | ".join(row) + " |")
        if rows:
            sep = "| " + " | ".join(["---"] * len(rows[0].split("|")[1:-1])) + " |"
            rows.insert(1, sep)
        return "\n".join(rows) or "[Empty CSV]"
    except Exception as e:
        return f"[ERROR reading CSV: {e}]"


def extract_drawio(path: Path) -> str:
    """
    draw.io files are XML. We extract:
    - Shape labels (the text on nodes)
    - Edge labels (connections between nodes)
    - A plain-English description of the diagram structure
    """
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        # draw.io files are sometimes compressed — check for mxGraphModel
        if "<mxGraphModel" not in raw:
            return f"[draw.io file appears compressed or empty — raw size: {len(raw)} bytes]\n\nRaw content:\n```\n{raw[:2000]}\n```"

        root = ET.fromstring(raw)
        nodes = {}   # id -> label
        edges = []   # (source_label, target_label, edge_label)

        for cell in root.iter("mxCell"):
            cell_id    = cell.get("id", "")
            value      = cell.get("value", "").strip()
            style      = cell.get("style", "")
            source     = cell.get("source", "")
            target     = cell.get("target", "")
            is_edge    = cell.get("edge", "0") == "1"
            is_vertex  = cell.get("vertex", "0") == "1"

            if is_vertex and value:
                # Strip HTML tags from labels
                clean_value = re.sub(r"<[^>]+>", " ", value).strip()
                nodes[cell_id] = clean_value

            if is_edge:
                edges.append((source, target, value))

        parts = []

        if nodes:
            parts.append("## Shapes / Nodes\n")
            for node_id, label in nodes.items():
                parts.append(f"- {label}")

        if edges:
            parts.append("\n## Connections\n")
            for src, tgt, label in edges:
                src_label = nodes.get(src, src or "?")
                tgt_label = nodes.get(tgt, tgt or "?")
                edge_label = f' ("{label}")' if label else ""
                parts.append(f"- {src_label} → {tgt_label}{edge_label}")

        if not parts:
            return "[draw.io file parsed but no shapes or connections found]"

        return "\n".join(parts)

    except ET.ParseError as e:
        return f"[ERROR parsing draw.io XML: {e}]\n\nRaw (first 1000 chars):\n```\n{raw[:1000]}\n```"
    except Exception as e:
        return f"[ERROR reading draw.io: {e}]"


def extract_image_vision(path: Path, anthropic_client) -> str:
    """
    Send image to Claude Vision for description.
    Returns a detailed description useful for notes.
    """
    try:
        from PIL import Image
        img = Image.open(path)
        width, height = img.size
        mode = img.mode
        size_kb = path.stat().st_size // 1024

        # Resize if very large to save API tokens
        max_dim = 1568
        if max(width, height) > max_dim:
            ratio = max_dim / max(width, height)
            img = img.resize((int(width * ratio), int(height * ratio)), Image.LANCZOS)

        # Convert to RGB if needed (for JPEG encoding)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        import io
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        img_b64 = base64.standard_b64encode(buf.getvalue()).decode("utf-8")

        response = anthropic_client.messages.create(
            model="claude-opus-4-5",
            max_tokens=1000,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/jpeg",
                            "data": img_b64
                        }
                    },
                    {
                        "type": "text",
                        "text": (
                            "This image is from a university student's school folder. "
                            "Describe what you see in detail, focusing on any text, diagrams, "
                            "charts, tables, formulas, or academic content visible. "
                            "If it's a diagram or flowchart, describe the structure and relationships. "
                            "If it contains text, transcribe the key parts. "
                            "Be thorough — this will be used as a study note."
                        )
                    }
                ]
            }]
        )

        description = response.content[0].text
        meta = f"> Image dimensions: {width}×{height}px | Size: {size_kb}KB | Mode: {mode}\n\n"
        return meta + description

    except ImportError:
        return "[ERROR: Pillow not installed — pip install Pillow]\n[Skipping vision analysis]"
    except Exception as e:
        return f"[ERROR analysing image with Claude Vision: {e}]"


def extract_code_or_text(path: Path) -> str:
    try:
        content = path.read_text(encoding="utf-8", errors="replace")
        ext = path.suffix.lower().lstrip(".")
        return f"```{ext}\n{content}\n```"
    except Exception as e:
        return f"[ERROR reading file: {e}]"


def extract_file(path: Path, use_vision: bool = True, anthropic_client=None) -> str:
    ext = path.suffix.lower()

    if ext == ".pdf":             return extract_pdf(path)
    if ext in (".docx", ".doc"): return extract_docx(path)
    if ext in SLIDE_EXTS:        return extract_pptx(path)
    if ext in (".xlsx", ".xls"): return extract_xlsx(path)
    if ext == ".csv":             return extract_csv(path)
    if ext == ".drawio":          return extract_drawio(path)
    if ext in CODE_EXTS:          return extract_code_or_text(path)
    if ext in (".txt", ".md"):    return extract_code_or_text(path)

    if ext in IMAGE_EXTS:
        if use_vision and anthropic_client:
            return extract_image_vision(path, anthropic_client)
        else:
            size_kb = path.stat().st_size // 1024
            return f"[Image file: {path.name} | {size_kb}KB]\n[Vision analysis skipped — run without --no-vision to analyse]"

    return f"[Unsupported file type: {ext}]"


def build_note(file_path: Path, content: str) -> str:
    try:
        rel = file_path.relative_to(SCHOOL_DIR)
    except ValueError:
        rel = Path(file_path.name)

    parts = list(rel.parts)
    subject  = parts[0] if len(parts) > 1 else "general"
    subfolder = " > ".join(parts[:-1]) if len(parts) > 2 else subject
    now = datetime.now().strftime("%Y-%m-%d")

    return f"""---
title: "{file_path.stem}"
source: "{file_path}"
subject: "{subject}"
path: "{subfolder}"
type: "{file_path.suffix.lower().lstrip('.')}"
scanned: "{now}"
tags: [education, {subject.lower().replace(' ', '-')}]
---

# {file_path.stem}

> **Source:** `{rel}`
> **Type:** {file_path.suffix.upper()} | **Subject:** {subject}

---

{content}
"""


# ── Main ───────────────────────────────────────────────────────────────────

def scan_and_convert(subject_filter: str = None, use_vision: bool = True):
    # Optionally init Anthropic client for vision
    anthropic_client = None
    if use_vision:
        try:
            import anthropic
            anthropic_client = anthropic.Anthropic()
            print("✅ Claude Vision enabled — images will be analysed")
        except ImportError:
            print("⚠️  anthropic package not found — pip install anthropic")
            print("   Continuing without vision (images will be skipped)\n")
            use_vision = False
        except Exception as e:
            print(f"⚠️  Could not init Anthropic client: {e}")
            print("   Make sure ANTHROPIC_API_KEY is set as an environment variable")
            print("   Continuing without vision\n")
            use_vision = False

    # Determine scan root
    scan_root = SCHOOL_DIR
    if subject_filter:
        for d in SCHOOL_DIR.rglob("*"):
            if d.is_dir() and subject_filter.lower() in d.name.lower():
                scan_root = d
                print(f"🎯 Filtering to: {d}")
                break

    print(f"\n{'='*60}")
    print(f"  School Folder Scanner")
    print(f"  Source : {scan_root}")
    print(f"  Output : {OUTPUT_DIR}")
    print(f"  Vision : {'ON' if use_vision else 'OFF'}")
    print(f"{'='*60}\n")

    processed = skipped = errors = 0

    for file_path in sorted(scan_root.rglob("*")):
        if not file_path.is_file():
            continue
        if any(skip in file_path.parts for skip in SKIP_DIRS):
            continue
        if file_path.name.lower() in SKIP_FILES:
            continue
        if file_path.name.startswith(".") or file_path.name.startswith("$"):
            continue
        if file_path.suffix.lower() not in ALL_SUPPORTED:
            skipped += 1
            continue

        # Mirror folder structure in output
        try:
            rel = file_path.relative_to(SCHOOL_DIR)
        except ValueError:
            rel = file_path.relative_to(scan_root)

        out_folder = OUTPUT_DIR / rel.parent
        out_folder.mkdir(parents=True, exist_ok=True)
        out_file = out_folder / (safe_filename(file_path.stem) + ".md")

        # Skip up-to-date files
        if out_file.exists():
            if out_file.stat().st_mtime >= file_path.stat().st_mtime:
                skipped += 1
                continue

        ext_label = file_path.suffix.upper().lstrip(".")
        print(f"  [{ext_label:6}] {rel}")

        try:
            content = extract_file(file_path, use_vision=use_vision, anthropic_client=anthropic_client)
            note    = build_note(file_path, content)
            out_file.write_text(note, encoding="utf-8")
            print(f"           ✅ → education/{rel.parent}/{out_file.name}")
            processed += 1
        except Exception as e:
            print(f"           ❌ Error: {e}")
            errors += 1

    print(f"\n{'='*60}")
    print(f"  Done!  Processed: {processed} | Skipped: {skipped} | Errors: {errors}")
    print(f"  Notes saved to: {OUTPUT_DIR}")
    print(f"{'='*60}\n")
    if processed > 0:
        print("💡 Claude can now read all notes directly via Obsidian MCP.")
        print('   Just say "load my [subject] context" in any session.\n')


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Scan school folder → Obsidian markdown notes")
    parser.add_argument("--subject",   type=str, default=None, help='Filter to one subject e.g. "Database Management"')
    parser.add_argument("--no-vision", action="store_true",    help="Skip Claude Vision image analysis (faster)")
    args = parser.parse_args()
    scan_and_convert(subject_filter=args.subject, use_vision=not args.no_vision)
